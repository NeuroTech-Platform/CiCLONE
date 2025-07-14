import os
import re
from datetime import datetime
from pathlib import Path
from typing import List, Tuple, Dict

from ciclone.utils.file_utils import FileUtils
from pptx import Presentation


class SchemaProcessor:
    """Service for processing schema files including PowerPoint conversion and multi-file handling."""
    
    @classmethod
    def is_supported_file(cls, file_path: str) -> bool:
        """Check if file is supported (image or PowerPoint)."""
        return (FileUtils.is_file_type(file_path, FileUtils.IMAGE_EXTENSIONS) or 
                FileUtils.is_file_type(file_path, FileUtils.POWERPOINT_EXTENSIONS))
    
    @classmethod
    def is_powerpoint_file(cls, file_path: str) -> bool:
        """Check if file is a PowerPoint file."""
        return FileUtils.is_file_type(file_path, FileUtils.POWERPOINT_EXTENSIONS)
    
    @classmethod
    def is_image_file(cls, file_path: str) -> bool:
        """Check if file is an image file."""
        return FileUtils.is_file_type(file_path, FileUtils.IMAGE_EXTENSIONS)
    
    @classmethod
    def convert_powerpoint_to_images(cls, ppt_path: str, output_dir: str, 
                                   output_format: str = 'png') -> Tuple[bool, List[str], str]:
        """
        Convert PowerPoint slides to structured content with images and text using python-pptx.
        
        Args:
            ppt_path: Path to PowerPoint file
            output_dir: Directory to save converted content
            output_format: Output image format ('png', 'jpg', 'tiff')
        
        Returns:
            Tuple of (success, list_of_output_paths, error_message)
        """
        try:
            # Ensure output directory exists
            if not FileUtils.ensure_directory(output_dir):
                return False, [], f"Cannot create output directory: {output_dir}"
            
            # Get base filename without extension
            base_name = Path(ppt_path).stem
            
            # Convert PowerPoint using python-pptx
            return cls._convert_ppt_with_pptx(ppt_path, output_dir, base_name)
                
        except Exception as e:
            return False, [], f"PowerPoint conversion failed: {str(e)}"
    
    @classmethod
    def _convert_ppt_with_pptx(cls, ppt_path: str, output_dir: str, base_name: str) -> Tuple[bool, List[str], str]:
        """Convert PowerPoint using python-pptx for both slide structure and images."""
        try:
            # Step 1: Use python-pptx to extract slide structure and text
            print("Extracting slide structure with python-pptx...")
            slide_contents = cls._extract_slides_with_pptx(ppt_path)
            
            # Step 2: Extract images directly from slides using python-pptx
            print("Extracting images with python-pptx...")
            slide_images = cls._extract_images_by_slide(ppt_path, output_dir, base_name)
            
            output_paths = []
            total_images = sum(len(images) for images in slide_images.values())
            
            # Add all image paths to output
            for slide_num, images in slide_images.items():
                for image_path in images:
                    if os.path.exists(image_path):
                        output_paths.append(image_path)
            
            # Step 3: Combine slide structure with images
            enhanced_markdown_content = cls._create_hybrid_markdown_with_slide_images(
                slide_contents, slide_images, base_name, ppt_path
            )
            
            # Save enhanced markdown file
            markdown_filename = f"{base_name}_slides.md"
            markdown_path = os.path.join(output_dir, markdown_filename)
            
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(enhanced_markdown_content)
            
            output_paths.insert(0, markdown_path)
            
            print(f"Created enhanced markdown: {markdown_filename}")
            return True, output_paths, f"Successfully converted PowerPoint using python-pptx: {len(slide_contents)} slides, {total_images} images extracted."
            
        except Exception as e:
            print(f"PowerPoint conversion failed: {str(e)}")
            return False, [], f"PowerPoint conversion failed: {str(e)}"
    
    @classmethod
    def _extract_slides_with_pptx(cls, ppt_path: str) -> List[Dict[str, str]]:
        """
        Extract slide content using python-pptx to preserve slide structure.
        
        Args:
            ppt_path: Path to PowerPoint file
        
        Returns:
            List of dictionaries with slide titles and content
        """
        try:
            prs = Presentation(ppt_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                title = ""
                content = ""
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # First meaningful text becomes the title
                        if not title and len(text) < 100:
                            title = text
                        else:
                            # All other text becomes content
                            if content:
                                content += "\n\n"
                            content += text
                
                # Clean up title and content
                if not title:
                    title = f"Slide {slide_num}"
                
                slides.append({
                    "title": title,
                    "content": content or ""
                })
                
                print(f"Extracted slide {slide_num}: {title[:50]}...")
            
            return slides
            
        except Exception as e:
            print(f"Warning: Could not extract slides with python-pptx: {str(e)}")
            return []
    
    @classmethod
    def _extract_images_by_slide(cls, ppt_path: str, output_dir: str, base_name: str) -> Dict[int, List[str]]:
        """
        Extract images from PowerPoint, organized by slide number.
        
        Args:
            ppt_path: Path to PowerPoint file
            output_dir: Directory to save images
            base_name: Base name for image files
        
        Returns:
            Dictionary mapping slide numbers to lists of image filenames
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image as PILImage
            import io
            
            prs = Presentation(ppt_path)
            slide_images = {}
            global_image_count = 1
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_image_list = []
                
                # Look for images in the slide
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # Extract image from shape
                            image_stream = io.BytesIO(shape.image.blob)
                            image = PILImage.open(image_stream)
                            
                            # Save image
                            image_filename = f"{base_name}_slide_{slide_num}_{len(slide_image_list)+1:02d}.png"
                            image_path = os.path.join(output_dir, image_filename)
                            
                            # Convert to RGB if necessary and save
                            if image.mode in ('RGBA', 'LA', 'P'):
                                rgb_image = PILImage.new('RGB', image.size, (255, 255, 255))
                                if image.mode == 'P':
                                    image = image.convert('RGBA')
                                rgb_image.paste(image, mask=image.split()[-1] if image.mode in ('RGBA', 'LA') else None)
                                rgb_image.save(image_path, 'PNG')
                            else:
                                image.save(image_path, 'PNG')
                            
                            slide_image_list.append(image_filename)
                            print(f"Extracted image from slide {slide_num}: {image_filename}")
                            global_image_count += 1
                            
                    except Exception as e:
                        print(f"Warning: Could not extract image from slide {slide_num}: {str(e)}")
                        continue
                
                if slide_image_list:
                    slide_images[slide_num] = slide_image_list
                    print(f"Slide {slide_num}: {len(slide_image_list)} images")
                else:
                    print(f"Slide {slide_num}: No images found")
            
            return slide_images
            
        except Exception as e:
            print(f"Warning: Could not extract images with python-pptx: {str(e)}")
            return {}
    
    @classmethod
    def _create_hybrid_markdown_with_slide_images(cls, slide_contents: List[Dict[str, str]], 
                                                 slide_images: Dict[int, List[str]], 
                                                 base_name: str, source_ppt_path: str) -> str:
        """
        Create markdown content with images properly matched to their original slides.
        
        Args:
            slide_contents: List of slide dictionaries from python-pptx
            slide_images: Dictionary mapping slide numbers to image filenames
            base_name: Base name for the presentation
            source_ppt_path: Path to the original PowerPoint file
        
        Returns:
            Enhanced markdown content with proper slide-image matching
        """
        total_images = sum(len(images) for images in slide_images.values())
        
        # Create header
        header = f"""# {base_name.replace('_', ' ').title()} - PowerPoint Conversion

Source File: {os.path.basename(source_ppt_path)}  
Converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  
Slides: {len(slide_contents)}
Images Extracted: {total_images}

---

"""
        
        # Create slides with their actual images
        formatted_slides = []
        
        for i, slide in enumerate(slide_contents):
            slide_num = i + 1
            
            # Create slide header
            formatted_slide = f"## Slide {slide_num}: {slide['title']}\n\n"
            
            # Add images that actually belong to this slide
            if slide_num in slide_images:
                for j, image_filename in enumerate(slide_images[slide_num]):
                    formatted_slide += f"![]({image_filename})\n"
            
            # Add slide content (only if there is content)
            if slide['content'].strip():
                formatted_slide += f"{slide['content']}\n\n"
            
            formatted_slides.append(formatted_slide)
        
        return header + ''.join(formatted_slides)


    @classmethod
    def copy_and_rename_images(cls, image_paths: List[str], output_dir: str, 
                              base_name: str = "schema") -> Tuple[bool, List[str], str]:
        """
        Copy multiple image files to output directory with consistent naming.
        
        Args:
            image_paths: List of source image file paths
            output_dir: Directory to copy images to
            base_name: Base name for renamed files
        
        Returns:
            Tuple of (success, list_of_copied_paths, error_message)
        """
        if not FileUtils.ensure_directory(output_dir):
            return False, [], f"Cannot create output directory: {output_dir}"
        
        copied_paths = []
        errors = []
        
        for i, image_path in enumerate(image_paths, 1):
            if not cls.is_image_file(image_path):
                continue
            
            # Generate new filename using utility
            file_ext = FileUtils.get_file_extension(image_path)
            new_filename = FileUtils.generate_filename(base_name, file_ext, i, len(image_paths))
            output_path = os.path.join(output_dir, new_filename)
            
            # Use centralized copy utility
            success, error = FileUtils.safe_copy_file(image_path, output_path)
            if success:
                copied_paths.append(output_path)
            else:
                errors.append(error)
        
        if copied_paths:
            error_msg = "; ".join(errors) if errors else ""
            return True, copied_paths, error_msg
        else:
            return False, [], "; ".join(errors) if errors else "No valid image files to copy"
    
    @classmethod
    def process_schema_files(cls, file_paths: List[str], output_dir: str, 
                           subject_name: str) -> Tuple[bool, List[str], str]:
        """
        Process schema files (images or PowerPoint) for a subject.
        
        Args:
            file_paths: List of file paths to process
            output_dir: Subject's output directory
            subject_name: Name of the subject (for naming)
        
        Returns:
            Tuple of (success, list_of_processed_files, error_message)
        """
        if not file_paths:
            return False, [], "No files provided"
        
        # Use the provided output directory directly (no subdirectory)
        if not FileUtils.ensure_directory(output_dir):
            return False, [], f"Cannot create output directory: {output_dir}"
        
        try:
            all_processed_files = []
            all_errors = []
            
            # Separate image files and PowerPoint files
            image_files = []
            powerpoint_files = []
            unsupported_files = []
            
            for file_path in file_paths:
                if not os.path.exists(file_path):
                    all_errors.append(f"File not found: {file_path}")
                    continue
                
                if cls.is_powerpoint_file(file_path):
                    powerpoint_files.append(file_path)
                elif cls.is_image_file(file_path):
                    image_files.append(file_path)
                else:
                    unsupported_files.append(file_path)
            
            # Process all image files together (for proper numbering)
            if image_files:
                base_name = f"{subject_name}_schema"
                success, copied_paths, error = cls.copy_and_rename_images(
                    image_files, output_dir, base_name
                )
                if success:
                    all_processed_files.extend(copied_paths)
                    print(f"Copied {len(copied_paths)} image files")
                if error:
                    all_errors.append(f"Image processing: {error}")
            
            # Process each PowerPoint file
            for ppt_file in powerpoint_files:
                success, output_paths, error = cls.convert_powerpoint_to_images(
                    ppt_file, output_dir, 'png'
                )
                if success:
                    all_processed_files.extend(output_paths)
                    if error:  # Warning message
                        all_errors.append(f"PowerPoint conversion warning: {error}")
                else:
                    all_errors.append(f"PowerPoint conversion failed: {error}")
            
            # Report unsupported files
            for unsupported_file in unsupported_files:
                all_errors.append(f"Unsupported file format: {unsupported_file}")
            
            # Combine results
            overall_success = len(all_processed_files) > 0
            error_message = "; ".join(all_errors) if all_errors else ""
            
            return overall_success, all_processed_files, error_message
            
        except Exception as e:
            return False, [], f"Schema processing failed: {str(e)}"
    
    @classmethod
    def get_supported_extensions_filter(cls) -> str:
        """Get file filter string for supported schema files."""
        file_types = {
            "Image files": FileUtils.IMAGE_EXTENSIONS,
            "PowerPoint files": FileUtils.POWERPOINT_EXTENSIONS
        }
        return FileUtils.create_file_filter(file_types) 